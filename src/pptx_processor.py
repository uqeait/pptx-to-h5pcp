import clr
import clr
import os
import sys
from pptx import Presentation
from PIL import Image

# Add reference to Spire.Presentation
clr.AddReference("Spire.Presentation")

from Spire.Presentation import Presentation as SpirePresentation

# Function to check if the file is a PowerPoint presentation
def is_pptx_file(filepath):
    return filepath.lower().endswith('.pptx')

# Function to open a PowerPoint file
def open_pptx(filepath):
    if not is_pptx_file(filepath):
        raise ValueError(f"The file {filepath} is not a PowerPoint (.pptx) file.")
    return Presentation(filepath)

# Function to extract hyperlinks from a shape
def extract_hyperlinks(shape):
    hyperlinks = []
    if shape.has_text_frame and shape.text:
        for paragraph in shape.text_frame.paragraphs:
            for run in paragraph.runs:
                if run.hyperlink:
                    hyperlinks.append({
                        "url": run.hyperlink.address,
                        "text": run.text,
                        "position": (shape.left, shape.top),
                        "width": shape.width,
                        "height": shape.height
                    })
    return hyperlinks

# Function to convert a slide to an image using Spire.Presentation
def convert_slide_to_image(slide, output_path, slide_index):
    # Create a new Spire.Presentation instance
    spire_presentation = SpirePresentation()

    # Add a blank slide (this will be replaced with the python-pptx slide)
    spire_slide = spire_presentation.slides.append()

    # TODO: Transfer content from python-pptx slide to Spire slide
    # This part is non-trivial and requires detailed implementation

    # Define the output image file path
    image_file = os.path.join(output_path, f"slide_{slide_index}.png")

    # Save the slide as an image
    spire_slide.save_as_image(image_file)

    # Dispose of the Spire.Presentation instance
    spire_presentation.dispose()

    return image_file

# Main function to process the PowerPoint file
def process_pptx_file(filepath, output_dir):
    presentation = open_pptx(filepath)
    slides_data = []

    for index, slide in enumerate(presentation.slides):
        slide_image_path = convert_slide_to_image(slide, output_dir, index)
        slide_content = {
            "image": slide_image_path,
            "hyperlinks": []
        }

        for shape in slide.shapes:
            hyperlinks = extract_hyperlinks(shape)
            slide_content["hyperlinks"].extend(hyperlinks)

        slides_data.append(slide_content)

    return slides_data

# Example usage
if __name__ == "__main__":
    pptx_file = "path_to_your_pptx_file.pptx"  # Replace with the actual file path
    output_dir = "path_to_output_directory"  # Replace with your desired output directory
    data = process_pptx_file(pptx_file, output_dir)
    print(data)  # For debugging purposes
